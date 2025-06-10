from django.shortcuts import render, redirect
from django.contrib import messages
from django.contrib.auth import logout
from django.contrib.auth.hashers import make_password, check_password
from pymongo import MongoClient
from datetime import datetime
import os, re
from django.core.files.storage import default_storage
from PyPDF2 import PdfReader
from sentence_transformers import SentenceTransformer
from django.http import JsonResponse
from django.views.decorators.http import require_POST, require_GET
from django.views.decorators.csrf import csrf_exempt
import requests
import json
import numpy as np
import fitz
import subprocess
import socket
import time
from docx import Document
from pptx import Presentation
from uuid import uuid4
import openpyxl
from bson import ObjectId
import calendar
import google.generativeai as genai

# MongoDB connection
client = MongoClient('mongodb://localhost:27017/')
db = client['eguru']
users_collection = db['users']
trainings_collection = db['trainings']
llm_collection = db['document']
collection = db['eguru-llm-index']
enrolled_courses_collection = db['enrolled_courses']
completed_trainings = db['completed_trainings']
materials_collection = db['materials']
review_collection = db['review']


CHUNK_SIZE = 500
CHUNK_OVERLAP = 50

embedder = SentenceTransformer('all-MiniLM-L6-v2')
genai.configure(api_key='AIzaSyAJGjyYHt6GX-dSUJgYLv3BUvWdS2nhfGg')
quiz_model = genai.GenerativeModel(model_name="models/gemini-1.5-flash")


def signin(request):
    if request.method == 'POST':
        pb_number = request.POST.get('pb_number')
        password = request.POST.get('password')
        user_data = users_collection.find_one({"pb_number": pb_number})

        if user_data and check_password(password, user_data.get('password')):
            request.session['pb_number'] = pb_number
            request.session['role'] = user_data.get('role').lower()

            if user_data.get('role').lower() == 'admin':
                return redirect('admin_dashboard')
            else:
                return redirect('user_dashboard')
        else:
            messages.error(request, 'Invalid PB Number or Password')

    return render(request, 'signin.html')

def signup(request):
    if request.method == 'POST':
        pb_number = request.POST.get('pb_number')
        password = request.POST.get('password')
        name = request.POST.get('name')
        gender = request.POST.get('gender')
        phone = request.POST.get('phone')
        email = request.POST.get('email')
        role = request.POST.get('role')
        division = request.POST.get('division')
        department = request.POST.get('department')
        designation = request.POST.get('designation')

        if users_collection.find_one({"pb_number": pb_number}):
            messages.error(request, 'PB Number already exists.')
        else:
            hashed_password = make_password(password)
            user_data = {
                "pb_number": pb_number,
                "password": hashed_password,
                "name": name,
                "gender": gender,
                "phone": phone,
                "email": email,
                "role": role,
                "division": division,
                "department": department,
                "designation": designation,
                "date_joined": datetime.now()
            }
            users_collection.insert_one(user_data)
            messages.success(request, 'Account created successfully. Please sign in.')
            return redirect('signin')

    return render(request, 'signup.html')

def change_password(request):
    if request.method == 'POST':
        pb_number = request.session.get('pb_number')
        old_password = request.POST.get('current_password')
        new_password = request.POST.get('new_password')

        user_data = users_collection.find_one({"pb_number": pb_number})

        if user_data and check_password(old_password, user_data.get('password')):
            hashed_new_password = make_password(new_password)
            users_collection.update_one(
                {"pb_number": pb_number},
                {"$set": {"password": hashed_new_password}}
            )
            messages.success(request, 'Password changed successfully.')
            return redirect('signin')
        else:
            messages.error(request, 'Invalid old password.')

    return render(request, 'signin.html')

def signout(request):
    # Clear specific session variables first
    if 'pb_number' in request.session:
        del request.session['pb_number']
    
    if 'role' in request.session:
        del request.session['role']
    logout(request)
    
    request.session.flush()
    
    request.session.clear_expired()
    
    messages.success(request, 'You have been logged out.')
    return redirect('signin')

def ensure_ollama_running():
    host = "localhost"
    port = 11434

    # Check if port is active (Ollama running)
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
        if sock.connect_ex((host, port)) == 0:
            print("✅ Ollama is already running.")
            return

    print("⏳ Starting Ollama with llama3 model...")
    try:
        subprocess.Popen(
            ["ollama", "run", "llama3"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        time.sleep(5)  # Wait briefly for server to start
    except Exception as e:
        print(f"❌ Failed to start Ollama: {e}")

def home(request):
    all_trainings = list(trainings_collection.find({}).limit(8))
    return render(request, 'home.html' , {'all_trainings': all_trainings})

def manage_trainings(request):
    if request.method == 'POST':
        title = request.POST.get('title')
        description = request.POST.get('description')
        category = request.POST.get('category')
        status = request.POST.get('status')
        instructor_pb_number = request.POST.get('instructor_pb_number')
        instructor_name = request.POST.get('instructor_name')
        start_date = request.POST.get('start_date')
        end_date = request.POST.get('end_date')
        image = request.FILES.get('image')

        if trainings_collection.find_one({"title": title, "instructor_pb_number": instructor_pb_number}):
            messages.error(request, 'Training with this title and instructor already exists.')
        elif not image:
            messages.error(request, 'Please upload an image.')
        else:
            # Validate image size (Max 2MB)
            if image.size > 2 * 1024 * 1024:
                messages.error(request, 'Image size exceeds 2MB.')
            else:
                # Validate image format
                valid_extensions = ['png', 'jpg', 'jpeg', 'webp']
                if image.name.split('.')[-1].lower() not in valid_extensions:
                    messages.error(request, 'Invalid image format.')
                else:
                    # Ensure the directory exists
                    image_directory = os.path.join('main','static' ,'training_images')
                    os.makedirs(image_directory, exist_ok=True)

                    # Save image to static folder
                    image_path = os.path.join(image_directory, image.name)
                    with open(image_path, 'wb+') as destination:
                        for chunk in image.chunks():
                            destination.write(chunk)

                    # Calculate duration
                    start_date_obj = datetime.strptime(start_date, '%Y-%m-%d')
                    end_date_obj = datetime.strptime(end_date, '%Y-%m-%d')
                    duration_days = (end_date_obj - start_date_obj).days

                    # Save to MongoDB
                    training_data = {
                        "title": title,
                        "description": description,
                        "category": category,
                        "status": status,
                        "instructor_pb_number": instructor_pb_number,
                        "instructor_name": instructor_name,
                        "start_date": start_date_obj,
                        "end_date": end_date_obj,
                        "duration": f"{duration_days} days",
                        "image_url": f"training_images\{image.name}"
                    }
                    trainings_collection.insert_one(training_data)
                    messages.success(request, 'Training added successfully.')
                    return redirect('admin_dashboard')

    return render(request, 'admin_dashboard.html')

def manage_llm(request):
    if request.method == "POST" and request.FILES.get("document"):
        document = request.FILES["document"]
        document_name = document.name.split('.')[0]

        if llm_collection.find_one({"name": document_name}):
            messages.error(request, 'Document with this name already exists.')
        elif not document:
            messages.error(request, 'Please upload a document.')
        else:
            # Validate file size (Max 100MB)
            if document.size > 100 * 1024 * 1024:
                messages.error(request, 'File size exceeds 100MB.')
            else:
                # Validate file type
                valid_extensions = ['pdf', 'ppt', 'pptx', 'doc', 'docx']
                if document.name.split('.')[-1].lower() not in valid_extensions:
                    messages.error(request, 'Invalid file format.')
                else:
                    # Ensure directory exists
                    document_directory = os.path.join('main','static','llm_documents')
                    os.makedirs(document_directory, exist_ok=True)

                    
                    document_path = os.path.join(document_directory, document.name)
                    with open(document_path, 'wb+') as destination:
                        for chunk in document.chunks():
                            destination.write(chunk)

                    # Save to MongoDB
                    document_data = {
                        "name": document_name,
                        "type": document.name.split('.')[-1].upper(),
                        "upload_date": datetime.now(),
                        "file_url": f"llm_documents\{document.name}"
                    }
                    llm_collection.insert_one(document_data)
                    # Index the document
                    index_single_document(document_path)
                    messages.success(request, 'Document uploaded successfully.')
                    return redirect('admin_dashboard')

    return render(request, 'admin_dashboard.html')

# Admin Dashboard with session validation
def admin_dashboard(request):
    
    pb_number = request.session.get('pb_number')
    
    user_data = users_collection.find_one({"pb_number": pb_number})
    all_users = list(users_collection.find({}))
    all_trainings = list(trainings_collection.find({}))
    enrollment = list(enrolled_courses_collection.find({}))
    completed = list(completed_trainings.find({}))    
    trainings_with_id = [
        {**training, 'training_id': str(training['_id'])} for training in all_trainings
    ]
    enroll_with_id = [
        {**enroll, 'training_id': str(enroll['training_id'])} for enroll in enrollment
    ]
    completed_with_id = [
        {**comp, 'training_id': str(comp['training_id'])} for comp in completed
    ]
    uploaded_documents = list(llm_collection.find({}))
    
    # ---------------- Month-wise Aggregation ----------------
    completions_month = list(completed_trainings.aggregate([
        {"$match": {"status": "completed"}},
        {"$group": {
            "_id": {"$dateToString": {"format": "%Y-%m", "date": "$completion_date"}},
            "count": {"$sum": 1}
        }},
        {"$sort": {"_id": 1}}
    ]))
    newusers_month = list(users_collection.aggregate([
        {"$group": {
            "_id": {"$dateToString": {"format": "%Y-%m", "date": "$date_joined"}},
            "count": {"$sum": 1}
        }},
        {"$sort": {"_id": 1}}
    ]))
    completions_year = list(completed_trainings.aggregate([
        {"$match": {"status": "completed",
                    "pb_number": pb_number}},
        {"$group": {
            "_id": {"$dateToString": {"format": "%Y", "date": "$completion_date"}},
            "count": {"$sum": 1}
        }},
        {"$sort": {"_id": 1}}
    ]))
    newusers_year = list(users_collection.aggregate([
        {"$group": {
            "_id": {"$dateToString": {"format": "%Y", "date": "$date_joined"}},
            "count": {"$sum": 1}
        }},
        {"$sort": {"_id": 1}}
    ]))
    months = [f"2025-{str(m).zfill(2)}" for m in range(1, 13)]
    years = sorted(list(set([entry['_id'] for entry in newusers_year + completions_year])))
    cm_dict = {e['_id']: e['count'] for e in completions_month}
    nm_dict = {c['_id']: c['count'] for c in newusers_month}
    cy_dict = {e['_id']: e['count'] for e in completions_year}
    ny_dict = {c['_id']: c['count'] for c in newusers_year}
    monthly_complete = [cm_dict.get(month, 0) for month in months]
    monthly_newusers = [nm_dict.get(month, 0) for month in months]
    yearly_complete = [cy_dict.get(year, 0) for year in years]
    yearly_newusers = [ny_dict.get(year, 0) for year in years]

    context ={
        'all_users': all_users,
        'admin': user_data,
        'trainings': trainings_with_id,
        'uploaded_documents': uploaded_documents,
        'enrollment': enroll_with_id,
        'completed': completed_with_id,
        'month_labels': [calendar.month_abbr[int(m[5:])] for m in months],
        'month_new_users': monthly_newusers,
        'month_complete': monthly_complete,
        'year_labels': years,
        'year_new_users': yearly_newusers,
        'year_complete': yearly_complete,
    }
    return render(request, 'admin_dashboard.html', context)

def user_dashboard(request):
    pb_number = request.session.get('pb_number')
    user_data = users_collection.find_one({"pb_number": pb_number})
    all_trainings = list(trainings_collection.find({}))
    uploaded_documents = list(llm_collection.find({}))
    trainings_with_id = [
        {**training, 'training_id': str(training['_id'])} for training in all_trainings
    ]
    enrollment = list(enrolled_courses_collection.find({"pb_number": pb_number}))
    completed = list(completed_trainings.find({"pb_number": pb_number}))    
    enroll_with_id = [
        {**enroll, 'training_id': str(enroll['training_id'])} for enroll in enrollment
    ]
    completed_with_id = [
        {**comp, 'training_id': str(comp['training_id'])} for comp in completed
    ]
    enrollments_month = list(enrolled_courses_collection.aggregate([
        {"$match": {"status": "enrolled",
                    "pb_number": pb_number}},
        {"$group": {
            "_id": {"$dateToString": {"format": "%Y-%m", "date": "$enrollment_date"}},
            "count": {"$sum": 1}
        }},
        {"$sort": {"_id": 1}}
    ]))
    completions_month = list(completed_trainings.aggregate([
        {"$match": {"status": "completed",
                    "pb_number": pb_number}},
        {"$group": {
            "_id": {"$dateToString": {"format": "%Y-%m", "date": "$completion_date"}},
            "count": {"$sum": 1}
        }},
        {"$sort": {"_id": 1}}
    ]))
    enrollments_year = list(enrolled_courses_collection.aggregate([
        {"$match": {"status": "enrolled",
                    "pb_number": pb_number}},
        {"$group": {
            "_id": {"$dateToString": {"format": "%Y", "date": "$enrollment_date"}},
            "count": {"$sum": 1}
        }},
        {"$sort": {"_id": 1}}
    ]))
    completions_year = list(completed_trainings.aggregate([
        {"$match": {"status": "completed",
                    "pb_number": pb_number}},
        {"$group": {
            "_id": {"$dateToString": {"format": "%Y", "date": "$completion_date"}},
            "count": {"$sum": 1}
        }},
        {"$sort": {"_id": 1}}
    ]))
    categories = ['aerodynamics', 'aviation', 'personality development', 'technical', 'corporate', 'management']
    category_stats = {cat: {"enrolled": 0, "completed": 0} for cat in categories}
    category_labels = [category.title() for category in categories]
    for enroll in enroll_with_id:
        category = enroll.get("category")
        if category in category_stats:
            category_stats[category]["enrolled"] += 1
    
    for comp in completed_with_id:
        category = comp.get("category")
        if category in category_stats:
            category_stats[category]["completed"] += 1
    enrolled_counts = [category_stats[cat]["enrolled"] for cat in categories]
    completed_counts = [category_stats[cat]["completed"] for cat in categories]
    
    months = [f"2025-{str(m).zfill(2)}" for m in range(1, 13)]
    years = sorted(list(set([entry['_id'] for entry in enrollments_year + completions_year])))
    
    # Convert to dict for lookup
    em_dict = {e['_id']: e['count'] for e in enrollments_month}
    cm_dict = {c['_id']: c['count'] for c in completions_month}
    ey_dict = {e['_id']: e['count'] for e in enrollments_year}
    cy_dict = {c['_id']: c['count'] for c in completions_year}

    # Fill zero where missing
    monthly_enroll = [em_dict.get(month, 0) for month in months]
    monthly_complete = [cm_dict.get(month, 0) for month in months]
    yearly_enroll = [ey_dict.get(year, 0) for year in years]
    yearly_complete = [cy_dict.get(year, 0) for year in years]
    context = {
        'user': user_data,
        'trainings': trainings_with_id,
        'uploaded_documents': uploaded_documents,
        'enrollment': enroll_with_id,
        'completed': completed_with_id,
        'month_labels': [calendar.month_abbr[int(m[5:])] for m in months],
        'month_enroll': monthly_enroll,
        'month_complete': monthly_complete,
        'year_labels': years,
        'year_enroll': yearly_enroll,
        'year_complete': yearly_complete,
        "category_labels": category_labels,
        "category_enrolled": enrolled_counts,
        "category_completed": completed_counts,
    }
    return render(request, 'user_dashboard.html', context)

def create_user(request):
    if request.method == 'POST':
        pb_number = request.POST.get('pb_number')
        password = request.POST.get('password')
        name = request.POST.get('name')
        phone = request.POST.get('phone')
        email = request.POST.get('email')
        gender = request.POST.get('gender')
        role = request.POST.get('role')
        division = request.POST.get('division')
        department = request.POST.get('department')
        designation = request.POST.get('designation')

        if users_collection.find_one({"pb_number": pb_number}):
            messages.error(request, 'PB Number already exists.')
        else:
            hashed_password = make_password(password)
            user_data = {
                "pb_number": pb_number,
                "password": hashed_password,
                "name": name,
                "phone": str(phone),
                "email": email,
                "gender": gender,
                "role": role,
                "division": division,
                "department": department,
                "designation": designation,
                "date_joined": datetime.now()
            }
            users_collection.insert_one(user_data)
            messages.success(request, 'User created successfully.')        
        return redirect('admin_dashboard')
    return render(request, 'admin_dashboard.html')

def import_users_excel(request):
    if request.method == 'POST' and request.FILES.get('excel_file'):
        excel_file = request.FILES['excel_file']
        wb = openpyxl.load_workbook(excel_file)
        sheet = wb.active

        for i, row in enumerate(sheet.iter_rows(min_row=2, values_only=True)):  
            try:
                pb_number, password, name, phone, email, gender, role, division, department, designation = row
                hashed_password = make_password(str(password))

                user_data = {
                    "pb_number": str(pb_number),
                    "password": hashed_password,
                    "name": name,
                    "phone": str(phone),
                    "email": email,
                    "gender": gender,
                    "role": role,
                    "division": division,
                    "department": department,
                    "designation": designation,
                    "date_joined": datetime.now()
                }

                existing = users_collection.find_one({"pb_number": str(pb_number)})
                if not existing:
                    users_collection.insert_one(user_data)
                else:
                    print(f"Skipping duplicate PB: {pb_number}")
            except Exception as e:
                print(f"Row {i+2} failed: {e}")

        messages.success(request, "Users imported successfully!")
        return redirect('admin_dashboard')  # or wherever your dashboard is

    messages.error(request, "No file uploaded.")
    return redirect('admin_dashboard')

def edit_user(request, user_id):
    # Check if authenticated user is admin
    pb_number = request.session.get('pb_number')
    user_data = users_collection.find_one({"pb_number": pb_number})
    
    if not user_data or user_data.get('role', '').lower() != 'admin':
        messages.error(request, 'Access denied. Admin access only.')
        return redirect('signin')
    
    # Implementation for editing users
    messages.info(request, 'Edit user functionality will be implemented soon.')
    return redirect('admin_dashboard')

def delete_user(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        pb_number = data.get('pb_number')

        if not pb_number:
            return JsonResponse({'error': 'PB number not provided'}, status=400)

        result = users_collection.delete_one({'pb_number': pb_number})

        if result.deleted_count == 1:
            return JsonResponse({'success': True})
        else:
            return JsonResponse({'success': False, 'message': 'User not found'})
    return redirect('admin_dashboard')


# Semantic Search
def find_relevant_chunks(query, top_k=5):
    query_emb = embedder.encode(query)
    all_docs = list(collection.find({}))
    for doc in all_docs:
        doc["score"] = np.dot(doc["embedding"], query_emb)
    sorted_docs = sorted(all_docs, key=lambda x: x["score"], reverse=True)
    return sorted_docs[:top_k]

# Call Ollama with context
def ask_ollama_with_context(context_chunks, user_query):
    ensure_ollama_running()
    context_text = "\n\n".join([
        f"[{doc['filename']} - Page {doc['page_number']}]\n{doc['text']}"
        for doc in context_chunks
    ])

    messages = [
        {"role": "system", "content": "You are a helpful assistant. Use the provided document context to answer the question."},
        {"role": "user", "content": f"Context:\n{context_text}"},
        {"role": "user", "content": user_query}
    ]

    response = requests.post(
        "http://localhost:11434/api/chat",
        json={"model": "llama3", "messages": messages, "stream": False}
    )

    # ✅ Add this check
    try:
        json_response = response.json()
        if "message" not in json_response:
            return f"⚠️ Ollama error: {json_response.get('error', 'Unknown response')}"
        return json_response["message"]["content"]
    except Exception as e:
        return f"❌ Failed to decode Ollama response: {str(e)}"


# Django view to handle chat request
@csrf_exempt
def ask_eguru(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            query = data.get("query", "").strip()

            if not query:
                return JsonResponse({"response": "Please provide a valid question."}, status=400)

            chunks = find_relevant_chunks(query)
            answer = ask_ollama_with_context(chunks, query)

            return JsonResponse({"response": answer})

        except Exception as e:
            return JsonResponse({"response": f"❌ Error: {str(e)}"}, status=500)

    return JsonResponse({"response": "Invalid request method"}, status=405)

def chunk_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

def extract_pdf(file_path):
    doc = fitz.open(file_path)
    contents = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text()
        if text.strip():
            contents.append({
                "filename": os.path.basename(file_path),
                "page_number": i,
                "text": text
            })
    return contents

def extract_docx(file_path):
    doc = Document(file_path)
    full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    return [{
        "filename": os.path.basename(file_path),
        "page_number": 1,
        "text": full_text
    }]

def extract_pptx(file_path):
    prs = Presentation(file_path)
    contents = []
    for i, slide in enumerate(prs.slides, start=1):
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        if text.strip():
            contents.append({
                "filename": os.path.basename(file_path),
                "page_number": i,
                "text": text
            })
    return contents

def index_single_document(file_path):
    ext = file_path.lower()
    if ext.endswith(".pdf"):
        docs = extract_pdf(file_path)
    elif ext.endswith(".docx"):
        docs = extract_docx(file_path)
    elif ext.endswith(".pptx") or ext.endswith(".ppt"):
        docs = extract_pptx(file_path)
    else:
        return  # unsupported

    for doc in docs:
        chunks = chunk_text(doc["text"])
        for i, chunk in enumerate(chunks):
            embedding = embedder.encode(chunk).tolist()
            collection.insert_one({
                "chunk_id": str(uuid4()),
                "filename": doc["filename"],
                "page_number": doc["page_number"],
                "chunk_index": i,
                "text": chunk,
                "embedding": embedding
            })

def enroll_training(request, training_id):
    if request.method == 'POST':
        pb_number = request.POST.get('pb_number') 
        name = request.POST.get('name')
        training_id = request.POST.get('training_id')
        
        # Check if the training exists
        training = trainings_collection.find_one({"_id": ObjectId(training_id)})
        
        if not training:
            messages.error(request, 'Training not found.')
            return redirect('user_dashboard')
        
        # Check if the user is already enrolled in the training
        existing_enrollment = enrolled_courses_collection.find_one({
            "pb_number": pb_number,
            "training_id": ObjectId(training_id)
        })

        if existing_enrollment:
            messages.info(request, 'You are already enrolled in this training.')
            return redirect('user_dashboard')

        # Enroll the user in the course
        enrolled_courses_collection.insert_one({
            "pb_number": pb_number,
            "name": name,
            "training_id": ObjectId(training_id),
            "category": training['category'],
            "trngtitle": training['title'],
            "status": "pending",
            "enrollment_date": datetime.now()
        })

        # After successful enrollment, redirect the user to the dashboard
        messages.success(request, 'You have successfully enrolled in the training.')
        return redirect('user_dashboard')

def complete_training(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        pb_number = data.get('pb_number')
        training_id = data.get('training_id')
        tcategory = data.get('tcategory')

        training = trainings_collection.find_one({"_id": ObjectId(training_id)})

        # Check if already completed
        completed_course = completed_trainings.find_one({
            "pb_number": pb_number,
            "training_id": ObjectId(training_id)
        })

        if completed_course:
            return JsonResponse({
                "success": False,
                "message": "You have already completed this training."
            })

        # Insert completion record
        completed_trainings.insert_one({
            "pb_number": pb_number,
            "training_id": ObjectId(training_id),
            "category": tcategory,
            "status": "completed",
            "completion_date": datetime.now()
        })

        return JsonResponse({
            "success": True,
            "message": "Training marked as completed."
        })

def training_details(request, title):
    if request.method == 'POST':
        pb_number = request.POST.get('pb_number')
        training_id = request.POST.get('training_id')
        user_data = users_collection.find_one({"pb_number": str(pb_number)})
        training = trainings_collection.find_one({"_id": ObjectId(training_id)})
        enrollment = enrolled_courses_collection.find_one({
        "pb_number": pb_number,
        "training_id": ObjectId(training_id)
        })
        status = enrollment.get('status') if enrollment else None
        materials = list(materials_collection.find({
            "training_id": ObjectId(training_id)
        }))
        materials_with_id = [
        {**material, 'material_id': str(material['_id'])} for material in materials
    ]
        completed = completed_trainings.find_one({
            "pb_number": pb_number,
            "training_id": ObjectId(training_id)
        })
        compstatus = completed.get('status') if completed else None
        reviews = list(review_collection.find({
            "training_id": ObjectId(training_id)
        }))
        review_summary = generate_review_summary(reviews)
        context = {
            'training': training,
            'user': user_data,
            'training_id': training_id,
            'pb_number': pb_number,
            'status': status,
            'materials': materials_with_id,
            'reviews': reviews,
            'compstatus': compstatus,
            'enrollment': enrollment,
            'completed': completed,
            'ai_summary': review_summary
        }
        return render(request, 'training_details.html', context)

def add_materials(request):
    if request.method == 'POST':
        training_id = request.POST.get('training_id')
        title = request.POST.get('title')
        material_type = request.POST.get('type')
        duration = request.POST.get('duration')
        
        # Validate required fields
        if not title or not material_type:
            messages.error(request, 'Title and material type are required.')
            return redirect('admin_dashboard')
        
        # Base directory for all materials
        material_directory = os.path.join('main', 'static', 'training_materials')
        os.makedirs(material_directory, exist_ok=True)
        
        material_data = {
            'title': title,
            'type': material_type,
            'training_id': ObjectId(training_id),
            'duration': duration,
            'upload_date': datetime.now(),
            'uploaded_by': request.session.get('pb_number'),
            'status': 'active'
        }
        
        try:
            # Handle different material types
            if material_type == 'file':
                file = request.FILES.get('file')
                if not file:
                    messages.error(request, 'Please upload a PDF file.')
                    return redirect('admin_dashboard')
                
                unique_filename = f"{file.name}"
                file_path = os.path.join(material_directory, unique_filename)
                
                with open(file_path, 'wb+') as destination:
                    for chunk in file.chunks():
                        destination.write(chunk)
                
                material_data['file_url'] = f"training_materials\{unique_filename}"
            
            elif material_type == 'video':
                video = request.FILES.get('video')
                if not video:
                    messages.error(request, 'Please upload a video file.')
                    return redirect('admin_dashboard')
                
                unique_filename = f"{video.name}"
                video_path = os.path.join(material_directory, unique_filename)
                
                with open(video_path, 'wb+') as destination:
                    for chunk in video.chunks():
                        destination.write(chunk)
                
                material_data['file_url'] = f"training_materials\{unique_filename}"
            
            elif material_type == 'link':
                link = request.POST.get('link')
                if not link:
                    messages.error(request, 'Please provide a valid URL.')
                    return redirect('admin_dashboard')
                material_data['external_url'] = link
            
            else:
                messages.error(request, 'Invalid material type selected.')
                return redirect('admin_dashboard')
            
            
            materials_collection.insert_one(material_data)
            # Update training with new material
            trainings_collection.update_one(
                {"_id": ObjectId(training_id)},
                {"$push": {"materials": material_data['_id']}}
            )
            
            messages.success(request, 'Material added successfully.')
        
        except Exception as e:
            messages.error(request, f'Error adding material: {str(e)}')
        
        return redirect('admin_dashboard')
    
    return render(request, 'training_details.html')

def delete_material(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        training_id = data.get('training_id')
        material_id = data.get('material_id')

        if not material_id:
            return JsonResponse({'error': 'Material ID not provided'}, status=400)

        result = materials_collection.delete_one({'_id': ObjectId(material_id)})

        if result.deleted_count == 1:
            trainings_collection.update_one(
                {'_id': ObjectId(training_id)},
                {'$pull': {'materials': ObjectId(material_id)}}
            )
            return JsonResponse({'success': True})
        else:
            return JsonResponse({'success': False, 'message': 'Material not found'})
    return redirect('training_details')

def delete_training(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        training_id = data.get('training_id')

        if not training_id:
            return JsonResponse({'error': 'Material ID not provided'}, status=400)

        result = trainings_collection.delete_one({'_id': ObjectId(training_id)})

    return redirect('admin_dashboard')

def extract_text_from_pdf(pdf_path):
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text
    return text

def generate_mcq_from_text(text, num_questions=10):
    prompt = f"""
    Generate {num_questions} multiple-choice questions (MCQs) based on the provided text. Follow these rules:  

    1. **Question Design:**  
        - Phrase each question as a **standalone fact or application** (e.g., avoid "What is the main focus of...?" or "According to the text...").  
        - Prioritize **specific concepts, relationships, or calculations** over broad summaries.  

    2. **Options & Answers:**  
        - Include **4 options (A-D)**, with **one correct answer** and **three plausible but incorrect distractors**.  
        - Ensure all options are **similar in length and complexity**.
        - Provide exactly 4 options (A-D) for each question
        - Only one correct answer per question
        - Incorrect options should be plausible but clearly wrong
        - Avoid "All of the above" or "None of the above" options
        - Keep all options roughly equal in length

    3. Strict Formatting:
        Question [number]: [Question text]
        A) [Option A]
        B) [Option B]
        C) [Option C]
        D) [Option D]
        Correct Answer: [Letter only - A/B/C/D]

    4. **Prohibited:**  
    - No references to "the text," "module," or "this section."  
    - No **all/none of the above** or **overlapping options**. 
    Example output:
    ```
    Question 1: What percentage of Earth’s atmosphere is nitrogen?  
    A) 21%  
    B) 78%  
    C) 50%  
    D) 1%  
    Correct Answer: B  
    Question 2: What is the primary function of the ozone layer?
    A) To provide oxygen
    B) To protect against UV radiation
    C) To regulate temperature
    D) To support life
    Correct Answer: B
    ```
    5. **Text Context:**
    - Use the provided text to generate questions.
    - Ensure questions are relevant to the content and concepts presented.
    6. **Output Format:**
    - Provide the questions in a single block of text, with each question starting with "Question [number]:"
    

    Text: {text}
    """
    response = quiz_model.generate_content(prompt)
    raw_mcqs = response.text

    mcqs = []
    blocks = re.split(r"Question \d+:", raw_mcqs)[1:]

    for block in blocks:
        lines = block.strip().split('\n')
        question = lines[0].strip()
        options = {}
        correct = ""
        for line in lines[1:]:
            if re.match(r'[A-D]\)', line):
                key = line[0]
                options[key] = line[3:].strip()
            elif "Correct Answer:" in line:
                correct = line.split(":")[1].strip()

        mcqs.append({
            "question": question,
            "options": [options.get("A"), options.get("B"), options.get("C"), options.get("D")],
            "correct": ["A", "B", "C", "D"].index(correct)
        })
    return mcqs

@csrf_exempt
def generate_quiz(request):
    if request.method == 'POST' and request.FILES.get('material_file'):
        uploaded_file = request.FILES['material_file']
        file_path = default_storage.save(f'main/static/temp/{uploaded_file.name}', uploaded_file)

        # Extract and generate MCQs
        abs_path = os.path.join(default_storage.location, file_path)
        pdf_text = extract_text_from_pdf(abs_path)
        mcqs = generate_mcq_from_text(pdf_text)

        return JsonResponse({'success': True, 'questions': mcqs})
    return JsonResponse({'success': False, 'message': 'Invalid request'})

@csrf_exempt
def save_quiz(request):
    if request.method == 'POST':
        try:
            quiz_data = json.loads(request.POST.get('quiz_data'))
            material_data = {
                'title': quiz_data['title'],
                'type': 'quiz',
                'training_id': ObjectId(quiz_data['training_id']),
                'duration': quiz_data['duration'],
                'uploaded_by': request.session.get('pb_number'),
                'status': 'active',
                'questions': quiz_data['questions']
            }
            materials_collection.insert_one(material_data)

            return redirect('training_details')    
        
        except Exception as e:
            return JsonResponse({'success': False, 'message': str(e)})
    return JsonResponse({'success': False, 'message': 'Invalid request method'})

def accept(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        pb_number = data.get('pb_number')
        training_id = data.get('training_id')

        enrolled_courses_collection.update_one(
            {"pb_number": pb_number, "training_id": ObjectId(training_id)},
            {"$set": {"status": "enrolled"}}
        )
        
        return JsonResponse({"success": True, "message": "Enrollment accepted."})

    return JsonResponse({"success": False, "message": "Invalid request."})

def reject(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        pb_number = data.get('pb_number')
        training_id = data.get('training_id')

        result = enrolled_courses_collection.delete_one({
            "pb_number": pb_number,
            "training_id": ObjectId(training_id)
        })

        if result.deleted_count > 0:
            return JsonResponse({"success": True, "message": "Enrollment deleted."})
        else:
            return JsonResponse({"success": False, "message": "No enrollment found to delete."})

    return JsonResponse({"success": False, "message": "Invalid request."})

@csrf_exempt
@require_POST
def submit_review(request):
    try:
        data = json.loads(request.body)
        
        # Validate required fields
        required_fields = ['training_id','training_title', 'user', 'pb_number', 'rating', 'comment']
        if not all(field in data for field in required_fields):
            return JsonResponse({'error': 'Missing required fields'}, status=400)
        
        try:
            rating = int(data['rating'])
            if rating < 1 or rating > 5:
                raise ValueError
        except (ValueError, TypeError):
            return JsonResponse({'error': 'Invalid rating value'}, status=400)
        
        review = {
            "training_id": ObjectId(data['training_id']),
            "training_title": data['training_title'],
            "user": data['user'],
            "pb_number": data['pb_number'],
            "rating": rating,
            "comment": data['comment'].strip(),
            "timestamp": datetime.now()
        }
        
        # Insert into MongoDB
        review_collection.insert_one(review)
        
        return JsonResponse({
            "message": "Review submitted successfully",
            "review": {
                "user": review['user'],
                "rating": review['rating'],
                "comment": review['comment'],
                "timestamp": review['timestamp'].isoformat()
            }
        })
        
    except json.JSONDecodeError:
        return JsonResponse({'error': 'Invalid JSON data'}, status=400)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)

def generate_review_summary(reviews):
    if not reviews:
        return "No reviews available"
    
    # Prepare review text for summarization
    review_text = "\n".join([
        f"Training Title: {review['training_title']} - Rating: {review['rating']}/5 - Comment: {review['comment']}"
        for review in reviews
    ])
    
    prompt = f"""
    You are a professional training evaluator with years of experience reviewing high-impact skill development programs. Analyze the following training reviews and provide a concise 4–5 line summary that highlights the overall training quality, key strengths, and areas for improvement.

    Avoid technical jargon. Keep the tone professional, suitable for a diverse audience that may include students, educators, institutional heads, professionals, and other stakeholders.

    Reviews:
    {review_text}

    Summary:
    """

    
    try:
        response = quiz_model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        print(f"Error generating review summary: {e}")
        return "Summary unavailable"